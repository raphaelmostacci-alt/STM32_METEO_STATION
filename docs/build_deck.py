# -*- coding: utf-8 -*-
"""
Builds the STM32 Meteo Station presentation:
- 6 slides (per project spec)
- 3 speakers, 2 slides each (badge on each slide + notes du presentateur)
- Palette Ocean Gradient (bleu marine + teal + accent ambre)

Run: python build_deck.py
Output: STM32_METEO_STATION_presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------
NAVY      = RGBColor(0x06, 0x5A, 0x82)   # primaire
TEAL      = RGBColor(0x1C, 0x72, 0x93)   # secondaire
MIDNIGHT  = RGBColor(0x21, 0x29, 0x5C)   # fonce
AMBER     = RGBColor(0xF5, 0x9E, 0x0B)   # accent (speaker tag)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF8, 0xFA, 0xFC)   # fond clair
CARD_BG   = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TXT  = RGBColor(0x0F, 0x17, 0x2A)
MUTED_TXT = RGBColor(0x64, 0x74, 0x8B)
EMERALD   = RGBColor(0x10, 0xB9, 0x81)   # pour les checks


# ---------------------------------------------------------------------------
# Setup
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def add_blank_slide():
    """Adds a slide using the 'Blank' layout."""
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_rounded(slide, x, y, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.12
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *,
             size=14, bold=False, color=DARK_TXT, align=PP_ALIGN.LEFT,
             valign=MSO_ANCHOR.TOP, font="Calibri", italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.text = ""
    if isinstance(text, str):
        lines = [text]
    else:
        lines = text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=DARK_TXT,
                font="Calibri", bullet_color=NAVY, gap=Pt(4)):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if i > 0:
            p.space_before = gap
        run_dot = p.add_run()
        run_dot.text = "●  "
        run_dot.font.name = font
        run_dot.font.size = Pt(size)
        run_dot.font.color.rgb = bullet_color
        run = p.add_run()
        run.text = item
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_speaker_chip(slide, speaker_name):
    """Amber chip top-right with the speaker name."""
    chip_w, chip_h = Inches(2.4), Inches(0.42)
    x = SLIDE_W - chip_w - Inches(0.45)
    y = Inches(0.35)
    chip = add_rounded(slide, x, y, chip_w, chip_h, AMBER)
    add_text(slide, x, y, chip_w, chip_h,
             "PARLE : " + speaker_name,
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_header(slide, number_label, title, speaker, *, dark=False):
    """Top header bar with slide number + title + speaker chip."""
    bar_color = NAVY if not dark else MIDNIGHT
    if dark:
        add_text(slide, Inches(0.6), Inches(0.4), Inches(4), Inches(0.4),
                 number_label, size=14, bold=True, color=AMBER, font="Consolas")
        add_text(slide, Inches(0.6), Inches(0.85), Inches(9), Inches(0.7),
                 title, size=30, bold=True, color=WHITE, font="Calibri")
    else:
        # white slide -> draw a navy bar across the top (1.5" tall)
        add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(1.4), bar_color)
        add_text(slide, Inches(0.6), Inches(0.25), Inches(5), Inches(0.4),
                 number_label, size=13, bold=True, color=AMBER, font="Consolas")
        add_text(slide, Inches(0.6), Inches(0.65), Inches(9.5), Inches(0.7),
                 title, size=24, bold=True, color=WHITE, font="Calibri")
    add_speaker_chip(slide, speaker)


def add_footer(slide, *, dark=False):
    color = WHITE if dark else MUTED_TXT
    add_text(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.35),
             "STM32 Meteo Station - NUCLEO-L152RE + X-NUCLEO-IKS01A3 - P23MRS 2026",
             size=10, color=color, font="Calibri", italic=True)
    add_text(slide, Inches(12.4), Inches(7.05), Inches(0.5), Inches(0.35),
             "", size=10, color=color)


def set_notes(slide, text):
    """Attach speaker notes."""
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    paras = text.strip().split("\n\n")
    for i, para in enumerate(paras):
        p = notes_tf.paragraphs[0] if i == 0 else notes_tf.add_paragraph()
        run = p.add_run()
        run.text = para
        run.font.size = Pt(12)
        run.font.name = "Calibri"


# ---------------------------------------------------------------------------
# Slide 1 - Titre / Application (Personne 1)
# ---------------------------------------------------------------------------
s1 = add_blank_slide()
set_bg(s1, MIDNIGHT)

# decor bands
add_rect(s1, Emu(0), Inches(6.0), SLIDE_W, Inches(0.15), AMBER)
add_rect(s1, Emu(0), Inches(6.15), SLIDE_W, Inches(0.10), TEAL)

# top-left small block "01"
add_text(s1, Inches(0.7), Inches(0.6), Inches(2), Inches(0.5),
         "01 / 06   APPLICATION", size=13, bold=True, color=AMBER, font="Consolas")
add_speaker_chip(s1, "Personne 1")

# big title
add_text(s1, Inches(0.7), Inches(1.6), Inches(11), Inches(1.4),
         "STATION METEO STM32", size=54, bold=True, color=WHITE, font="Calibri")
add_text(s1, Inches(0.7), Inches(2.8), Inches(11), Inches(0.7),
         "Acquisition environnementale temps reel", size=24, color=AMBER, italic=True, font="Calibri")

# divider
add_rect(s1, Inches(0.7), Inches(3.7), Inches(2), Inches(0.04), AMBER)

# bullets
add_text(s1, Inches(0.7), Inches(3.95), Inches(11), Inches(0.35),
         "Mesure de temperature, humidite et pression d'une piece, exposees en direct sur PC.",
         size=16, color=WHITE, font="Calibri")

# spec card bottom-left
card_x, card_y, card_w, card_h = Inches(0.7), Inches(4.6), Inches(5.7), Inches(1.3)
add_rounded(s1, card_x, card_y, card_w, card_h, NAVY)
add_text(s1, card_x + Inches(0.3), card_y + Inches(0.15), card_w - Inches(0.6), Inches(0.4),
         "PLATEFORME", size=11, bold=True, color=AMBER, font="Consolas")
add_text(s1, card_x + Inches(0.3), card_y + Inches(0.5), card_w - Inches(0.6), Inches(0.8),
         "NUCLEO-L152RE  +  X-NUCLEO-IKS01A3", size=18, bold=True, color=WHITE)
add_text(s1, card_x + Inches(0.3), card_y + Inches(0.88), card_w - Inches(0.6), Inches(0.4),
         "STM32L152RET6 - Cortex-M3 - 32 MHz", size=13, color=RGBColor(0xCB, 0xD5, 0xE1))

# spec card bottom-right
card2_x = Inches(7.0)
add_rounded(s1, card2_x, card_y, card_w, card_h, TEAL)
add_text(s1, card2_x + Inches(0.3), card_y + Inches(0.15), card_w - Inches(0.6), Inches(0.4),
         "CONTEXTE", size=11, bold=True, color=AMBER, font="Consolas")
add_text(s1, card2_x + Inches(0.3), card_y + Inches(0.5), card_w - Inches(0.6), Inches(0.8),
         "Projet de fin de module - P23MRS", size=18, bold=True, color=WHITE)
add_text(s1, card2_x + Inches(0.3), card_y + Inches(0.88), card_w - Inches(0.6), Inches(0.4),
         "ISEN - Promotion 2026", size=13, color=RGBColor(0xCB, 0xD5, 0xE1))

# footer date
add_text(s1, Inches(0.7), Inches(6.6), Inches(8), Inches(0.35),
         "Demonstration en direct - 20 minutes", size=12, color=WHITE, italic=True)

set_notes(s1, """[Personne 1] Bonjour a toutes et a tous.

Aujourd'hui on vous presente notre projet de fin de module: une station meteo embarquee sur STM32.

L'idee est simple: utiliser la carte NUCLEO-L152RE et le shield X-NUCLEO-IKS01A3 pour mesurer en continu la temperature, l'humidite et la pression d'une piece, et publier ces valeurs sur un PC via la liaison serie.

L'utilisateur peut meme entrer le nom de sa ville sur le terminal pour le voir apparaitre dans les logs - on en reparle plus loin.

Notre presentation va suivre 6 etapes : application, solution technique, defis, resolution, demo en direct, et conclusion. On parle a tour de role: moi pour les deux premieres slides, [Personne 2] pour les defis et la resolution, et [Personne 3] pour la demo et la conclusion.
""")


# ---------------------------------------------------------------------------
# Slide 2 - Solution technique (Personne 1)
# ---------------------------------------------------------------------------
s2 = add_blank_slide()
set_bg(s2, LIGHT_BG)
add_header(s2, "02 / 06   SOLUTION", "Architecture materielle et logicielle", "Personne 1")

# left card: HARDWARE
hw_x, hw_y, hw_w, hw_h = Inches(0.6), Inches(1.6), Inches(6.0), Inches(4.2)
add_rounded(s2, hw_x, hw_y, hw_w, hw_h, CARD_BG)
add_rect(s2, hw_x, hw_y, Inches(0.1), hw_h, NAVY)  # accent strip
add_text(s2, hw_x + Inches(0.35), hw_y + Inches(0.2), hw_w - Inches(0.55), Inches(0.4),
         "MATERIEL", size=12, bold=True, color=NAVY, font="Consolas")
add_text(s2, hw_x + Inches(0.35), hw_y + Inches(0.55), hw_w - Inches(0.55), Inches(0.5),
         "Une carte, un shield, trois capteurs", size=20, bold=True, color=DARK_TXT)

hw_items = [
    ("HTS221",   "Humidite + temperature - I2C"),
    ("LPS22HH",  "Pression atmospherique - I2C"),
    ("STTS751",  "Temperature de precision - I2C"),
    ("USART2",   "Terminal serie via ST-LINK - 115200 baud"),
    ("Bouton B1","Re-saisie ville par interruption EXTI13"),
]
row_y = hw_y + Inches(1.25)
for tag, desc in hw_items:
    add_text(s2, hw_x + Inches(0.35), row_y, Inches(1.6), Inches(0.35),
             tag, size=13, bold=True, color=TEAL, font="Consolas")
    add_text(s2, hw_x + Inches(2.0), row_y, hw_w - Inches(2.2), Inches(0.35),
             desc, size=13, color=DARK_TXT)
    row_y += Inches(0.45)

# right card: SOFTWARE
sw_x, sw_y, sw_w, sw_h = Inches(6.8), Inches(1.6), Inches(6.0), Inches(4.2)
add_rounded(s2, sw_x, sw_y, sw_w, sw_h, CARD_BG)
add_rect(s2, sw_x, sw_y, Inches(0.1), sw_h, TEAL)
add_text(s2, sw_x + Inches(0.35), sw_y + Inches(0.2), sw_w - Inches(0.55), Inches(0.4),
         "LOGICIEL", size=12, bold=True, color=TEAL, font="Consolas")
add_text(s2, sw_x + Inches(0.35), sw_y + Inches(0.55), sw_w - Inches(0.55), Inches(0.5),
         "Architecture flag-driven HAL C", size=20, bold=True, color=DARK_TXT)

sw_items = [
    ("CubeMX",       "Generation HAL + BSP IKS01A3 v11.1"),
    ("TIM6 IRQ",     "Tick d'echantillonnage a 1 Hz"),
    ("EXTI13 IRQ",   "Bouton B1 anti-rebond 50 ms"),
    ("USART2 RX IRQ","Reception caractere par caractere"),
    ("weather.c",    "Module separe : capteurs, LED, UART"),
]
row_y = sw_y + Inches(1.25)
for tag, desc in sw_items:
    add_text(s2, sw_x + Inches(0.35), row_y, Inches(2.0), Inches(0.35),
             tag, size=13, bold=True, color=NAVY, font="Consolas")
    add_text(s2, sw_x + Inches(2.4), row_y, sw_w - Inches(2.6), Inches(0.35),
             desc, size=13, color=DARK_TXT)
    row_y += Inches(0.45)

# bottom stats strip
stats_y = Inches(6.0)
stats = [
    ("3",     "capteurs"),
    ("3",     "IRQ"),
    ("1 Hz",  "echantillonnage"),
    ("115200","baud UART"),
]
stat_w = Inches(2.9)
gap = Inches(0.25)
total = stat_w * 4 + gap * 3
start_x = (SLIDE_W - total) // 2
for i, (big, small) in enumerate(stats):
    x = start_x + i * (stat_w + gap)
    add_rounded(s2, x, stats_y, stat_w, Inches(0.85), MIDNIGHT)
    add_text(s2, x, stats_y + Inches(0.08), stat_w, Inches(0.5),
             big, size=22, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    add_text(s2, x, stats_y + Inches(0.5), stat_w, Inches(0.35),
             small, size=12, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")

add_footer(s2)
set_notes(s2, """[Personne 1] Concretement, notre montage tient sur deux cartes empilees.

Cote materiel, la NUCLEO porte le MCU STM32L152, et le shield IKS01A3 apporte trois capteurs environnementaux : HTS221 pour l'humidite, LPS22HH pour la pression, et STTS751 pour la temperature. Les trois capteurs partagent le meme bus I2C1 sur les broches PB8 et PB9. La liaison UART passe par l'USART2 a 115200 bauds, directement par le port USB ST-LINK - donc pas besoin de cable supplementaire.

Cote logiciel, on a fait genere le squelette par CubeMX, et on a ajoute notre propre module weather.c qui regroupe l'init des capteurs, la lecture, le pilotage des LEDs et le format UART. Le projet repose entierement sur un pattern flag-driven : les ISR ne font que lever des drapeaux volatile, et tout le travail reel est fait dans la boucle while(1). Ca garantit qu'on ne bloque jamais le CPU dans une interruption.

Les chiffres cles en bas resument : trois capteurs, trois peripheriques en interruption - donc largement au-dessus des deux exigees - un cycle de mesure par seconde, et un debit UART de 115200 bauds.

Je passe la main a [Personne 2] qui va vous expliquer les defis qu'on a rencontres.
""")


# ---------------------------------------------------------------------------
# Slide 3 - Defi / problemes (Personne 2)
# ---------------------------------------------------------------------------
s3 = add_blank_slide()
set_bg(s3, LIGHT_BG)
add_header(s3, "03 / 06   DEFIS", "Quatre problemes a regler", "Personne 2")

# 2x2 grid of challenge cards
def challenge_card(slide, x, y, w, h, num, title, body):
    add_rounded(slide, x, y, w, h, CARD_BG)
    add_rect(slide, x, y, Inches(0.1), h, NAVY)
    # numero dans un cercle ambre
    circle_d = Inches(0.7)
    cx = x + Inches(0.35)
    cy = y + Inches(0.3)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, circle_d, circle_d)
    circle.fill.solid(); circle.fill.fore_color.rgb = AMBER
    circle.line.fill.background()
    circle.shadow.inherit = False
    add_text(slide, cx, cy, circle_d, circle_d, num,
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(1.2), y + Inches(0.35), w - Inches(1.4), Inches(0.5),
             title, size=18, bold=True, color=DARK_TXT)
    add_text(slide, x + Inches(0.35), y + Inches(1.2), w - Inches(0.6), h - Inches(1.3),
             body, size=13, color=MUTED_TXT)

cw, ch = Inches(6.1), Inches(2.35)
gx, gy = Inches(0.6), Inches(1.55)
gap_x, gap_y = Inches(0.15), Inches(0.2)

challenge_card(s3, gx,             gy,             cw, ch,
               "1", "Acquisition periodique",
               "Echantillonner pile chaque seconde sans bloquer le CPU pendant qu'une lecture I2C est lente. Il faut un cadencement materiel propre, pas un delay logiciel.")

challenge_card(s3, gx + cw + gap_x, gy,             cw, ch,
               "2", "Coexistence des IRQ",
               "Deux interruptions imposees par le sujet, on en a trois (TIM6, EXTI13, USART2 RX). Il faut les rendre coherentes : priorites, anti-rebond du bouton, callbacks non bloquants.")

challenge_card(s3, gx,             gy + ch + gap_y, cw, ch,
               "3", "Erreurs multi-capteurs",
               "Trois capteurs sur le meme bus I2C. Si l'un d'eux refuse de repondre, on doit continuer a publier les autres mesures et signaler l'erreur via une LED.")

challenge_card(s3, gx + cw + gap_x, gy + ch + gap_y, cw, ch,
               "4", "Interactivite UART",
               "Faire saisir la ville sans figer la mesure. Lire l'UART en interruption caractere par caractere et reprendre les mesures des que la ligne est validee.")

add_footer(s3)
set_notes(s3, """[Personne 2] Merci [Personne 1]. Concevoir une station meteo embarquee, ce n'est pas juste lire un capteur et l'afficher. On a identifie quatre defis principaux.

Premier defi : l'acquisition periodique. Il faut une mesure pile chaque seconde, sans utiliser de delay bloquant qui figerait tout le systeme. La solution evidente c'est un timer materiel.

Deuxieme defi : les interruptions. Le sujet impose qu'au moins deux peripheriques tournent en mode interruption. Nous, on en a trois - TIM6, le bouton EXTI13 et la reception UART - et il faut qu'ils cohabitent sans se marcher dessus. Notamment le bouton, qui rebondit physiquement, il a fallu un anti-rebond logiciel.

Troisieme defi : la robustesse. Trois capteurs sur le meme bus I2C, ca veut dire que si l'un tombe en panne ou est mal branche, il faut que le firmware continue a publier les deux autres et signale la panne clairement - dans notre cas par une LED rouge.

Quatrieme defi, celui qui a demande le plus de soin : permettre a l'utilisateur de saisir la ville sur le terminal. Concretement, recevoir l'UART caractere par caractere en interruption, mettre la telemetrie en pause pour laisser l'utilisateur taper, et reprendre des qu'il appuie sur Entree.

Maintenant je vais vous expliquer comment on a resolu tout ca.
""")


# ---------------------------------------------------------------------------
# Slide 4 - Resolution (Personne 2)
# ---------------------------------------------------------------------------
s4 = add_blank_slide()
set_bg(s4, LIGHT_BG)
add_header(s4, "04 / 06   RESOLUTION", "Pattern flag-driven et separation des couches", "Personne 2")

# left: architecture diagram
diag_x, diag_y = Inches(0.6), Inches(1.55)
diag_w, diag_h = Inches(6.5), Inches(5.2)
add_rounded(s4, diag_x, diag_y, diag_w, diag_h, CARD_BG)
add_rect(s4, diag_x, diag_y, Inches(0.1), diag_h, TEAL)
add_text(s4, diag_x + Inches(0.35), diag_y + Inches(0.2), diag_w - Inches(0.55), Inches(0.4),
         "ARCHITECTURE", size=12, bold=True, color=TEAL, font="Consolas")
add_text(s4, diag_x + Inches(0.35), diag_y + Inches(0.55), diag_w - Inches(0.55), Inches(0.5),
         "ISR ---> drapeau ---> main loop", size=19, bold=True, color=DARK_TXT)

# Three columns: ISR (left), main (center), OUT (right), arrows between.
# Inner zone of the card: x = 0.95 -> 6.75 (5.8" wide).
isr_w   = Inches(1.8)
main_w  = Inches(1.6)
out_w   = Inches(1.6)
arrow_w = Inches(0.4)

isr_x   = diag_x + Inches(0.35)
arrow1_x = isr_x + isr_w
main_x  = arrow1_x + arrow_w
arrow2_x = main_x + main_w
out_x   = arrow2_x + arrow_w

row_h = Inches(0.55)
gap = Inches(0.18)
y0 = diag_y + Inches(1.5)

irqs = [
    ("TIM6 IRQ",       "g_sample_flag = 1"),
    ("EXTI13 IRQ",     "g_ask_city = 1"),
    ("USART2 RX IRQ",  "g_city_ready = 1"),
]
for i, (label, var) in enumerate(irqs):
    y = y0 + i * (row_h + gap)
    add_rounded(s4, isr_x, y, isr_w, row_h, NAVY)
    add_text(s4, isr_x, y, isr_w, row_h,
             label, size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, font="Consolas")
    # arrow ISR -> main
    add_text(s4, arrow1_x, y, arrow_w, row_h, ">",
             size=22, bold=True, color=AMBER,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, font="Consolas")

# main loop box (centered)
main_h = row_h * 3 + gap * 2
add_rounded(s4, main_x, y0, main_w, main_h, AMBER)
add_text(s4, main_x, y0, main_w, main_h,
         "while (1)\n\nsensors I2C\nUART printf\nLEDs", size=11, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, font="Consolas")

outs = [
    ("UART line",       "METEO Ville=..."),
    ("LED status",      "L0 / L1"),
    ("Heartbeat",       "LD2 toggles"),
]
for i, (label, val) in enumerate(outs):
    y = y0 + i * (row_h + gap)
    # arrow main -> out
    add_text(s4, arrow2_x, y, arrow_w, row_h, ">",
             size=22, bold=True, color=AMBER,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, font="Consolas")
    add_rounded(s4, out_x, y, out_w, row_h, TEAL)
    add_text(s4, out_x, y, out_w, row_h,
             label, size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, font="Consolas")

# right: 4 small benefit cards
benefit_x = Inches(7.5)
benefit_y = Inches(1.55)
benefit_w = Inches(5.3)
benefit_h = Inches(1.18)
benefits = [
    ("Non bloquant",
     "Les ISR ne font que lever un drapeau et reviennent en < 1 us."),
    ("Anti-rebond",
     "HAL_GetTick() + fenetre 50 ms sur le bouton B1."),
    ("Pause interactive",
     "g_paused = 1 stoppe la telemetrie pendant la saisie ville."),
    ("Robustesse capteurs",
     "Init et lecture verifient BSP_ERROR_NONE, LED L1 sur faute."),
]
for i, (title, body) in enumerate(benefits):
    y = benefit_y + i * (benefit_h + Inches(0.12))
    add_rounded(s4, benefit_x, y, benefit_w, benefit_h, CARD_BG)
    add_rect(s4, benefit_x, y, Inches(0.08), benefit_h, AMBER)
    add_text(s4, benefit_x + Inches(0.25), y + Inches(0.15),
             benefit_w - Inches(0.4), Inches(0.4),
             title, size=15, bold=True, color=NAVY)
    add_text(s4, benefit_x + Inches(0.25), y + Inches(0.55),
             benefit_w - Inches(0.4), benefit_h - Inches(0.65),
             body, size=12, color=MUTED_TXT)

add_footer(s4)
set_notes(s4, """[Personne 2] Notre reponse, c'est une architecture flag-driven, illustree dans le diagramme a gauche.

A gauche les trois interruptions : TIM6, le bouton et la reception UART. Chacune ne fait qu'une chose - poser un drapeau volatile - et rend la main. Aucune lecture I2C, aucun printf dans une ISR.

Au centre, la boucle while(1) tourne en permanence et consulte ces drapeaux. Quand sample_flag est leve, elle lance la lecture I2C des trois capteurs, formate la ligne METEO et l'envoie sur l'UART. Quand ask_city est leve, elle affiche le prompt et met le drapeau g_paused a 1, ce qui suspend la telemetrie. Quand city_ready est leve, elle copie le buffer, affiche la confirmation et relance la mesure.

A droite, les benefices : tout est non bloquant, on a un anti-rebond robuste sur le bouton, l'interaction utilisateur ne se fait pas au detriment de la mesure, et la robustesse des capteurs est verifiee a chaque appel BSP.

[Personne 3], a toi pour la demo en direct.
""")


# ---------------------------------------------------------------------------
# Slide 5 - Demonstration (Personne 3)
# ---------------------------------------------------------------------------
s5 = add_blank_slide()
set_bg(s5, LIGHT_BG)
add_header(s5, "05 / 06   DEMO", "Demonstration en direct", "Personne 3")

# left: numbered steps
steps_x, steps_y = Inches(0.6), Inches(1.55)
steps_w, steps_h = Inches(6.3), Inches(5.2)
add_rounded(s5, steps_x, steps_y, steps_w, steps_h, CARD_BG)
add_rect(s5, steps_x, steps_y, Inches(0.1), steps_h, NAVY)
add_text(s5, steps_x + Inches(0.35), steps_y + Inches(0.2),
         steps_w - Inches(0.55), Inches(0.4),
         "SCENARIO", size=12, bold=True, color=NAVY, font="Consolas")
add_text(s5, steps_x + Inches(0.35), steps_y + Inches(0.55),
         steps_w - Inches(0.55), Inches(0.5),
         "Cinq etapes en direct", size=20, bold=True, color=DARK_TXT)

steps = [
    ("Reset",       "Banniere de boot + 3 capteurs prets."),
    ("Saisie ville","Tape 'Marseille' sur le terminal, validation Entree."),
    ("Mesure live", "Une ligne METEO par seconde, heartbeat LD2 a 1 Hz."),
    ("Variation",   "Souffler sur le shield - humidite + temperature montent."),
    ("Changer",     "Appui B1 -> nouvelle saisie ('Lyon'), la ville change."),
]
sy = steps_y + Inches(1.2)
for i, (label, desc) in enumerate(steps):
    add_text(s5, steps_x + Inches(0.35), sy, Inches(0.5), Inches(0.45),
             str(i+1), size=22, bold=True, color=AMBER, font="Consolas")
    add_text(s5, steps_x + Inches(0.95), sy, Inches(1.7), Inches(0.45),
             label, size=14, bold=True, color=NAVY, font="Consolas")
    add_text(s5, steps_x + Inches(2.7), sy, steps_w - Inches(2.9), Inches(0.7),
             desc, size=13, color=DARK_TXT)
    sy += Inches(0.78)

# right: terminal mockup
term_x, term_y = Inches(7.2), Inches(1.55)
term_w, term_h = Inches(5.55), Inches(5.2)
add_rounded(s5, term_x, term_y, term_w, term_h, MIDNIGHT)
add_text(s5, term_x + Inches(0.3), term_y + Inches(0.2),
         term_w - Inches(0.6), Inches(0.4),
         "TERMINAL  -  COM3  -  115200 8N1", size=11, bold=True, color=AMBER, font="Consolas")
# fake terminal lines
lines = [
    ("=== STM32 METEO STATION boot ===",                       WHITE),
    ("Board: NUCLEO-L152RE + X-NUCLEO-IKS01A3",                RGBColor(0xCB, 0xD5, 0xE1)),
    ("Capteurs HTS221 + LPS22HH + STTS751 prets.",             EMERALD),
    ("",                                                       WHITE),
    ("Entrez la ville : Marseille",                            AMBER),
    ("Ville selectionnee : Marseille",                         EMERALD),
    ("METEO Ville=Marseille T=22.45 C RH=47.80 % P=1015.2 hPa",WHITE),
    ("METEO Ville=Marseille T=22.47 C RH=47.92 % P=1015.2 hPa",WHITE),
    ("METEO Ville=Marseille T=22.51 C RH=48.10 % P=1015.3 hPa",WHITE),
    ("",                                                       WHITE),
    ("Entrez la ville : Lyon",                                 AMBER),
    ("Ville selectionnee : Lyon",                              EMERALD),
    ("METEO Ville=Lyon T=22.53 C RH=48.15 % P=1015.3 hPa",     WHITE),
]
ty = term_y + Inches(0.65)
for txt, col in lines:
    add_text(s5, term_x + Inches(0.3), ty, term_w - Inches(0.6), Inches(0.28),
             txt if txt else " ", size=11, color=col, font="Consolas")
    ty += Inches(0.30)

add_footer(s5)
set_notes(s5, """[Personne 3] Merci. Je vais vous montrer le projet en marche.

(Brancher la carte, ouvrir TeraTerm sur 115200 baud, appuyer sur reset.)

Etape 1, le reset. La banniere apparait, on voit que les trois capteurs - HTS221, LPS22HH et STTS751 - se sont initialises correctement.

Etape 2, je tape une ville. Notez qu'aucune ligne METEO n'apparait pendant que je tape - c'est la pause dont parlait [Personne 2]. Je valide avec Entree.

Etape 3, des que c'est valide, la telemetrie redemarre. Une ligne par seconde, avec la temperature, l'humidite et la pression mesurees, et le nom de la ville. La LED verte LD2 clignote au rythme du timer TIM6.

Etape 4, je souffle doucement sur le shield. Vous voyez que l'humidite monte de quelques pourcents, et la temperature aussi un peu - les capteurs reagissent bien au monde reel.

Etape 5, je veux changer de ville. J'appuie sur le bouton bleu B1 de la Nucleo - regardez, la telemetrie se met en pause, le prompt revient, et je peux taper une autre ville.

Voila, tout marche. On va terminer par la conclusion.
""")


# ---------------------------------------------------------------------------
# Slide 6 - Conclusion (Personne 3)
# ---------------------------------------------------------------------------
s6 = add_blank_slide()
set_bg(s6, MIDNIGHT)

# decor strip
add_rect(s6, Emu(0), Inches(0.0), SLIDE_W, Inches(0.15), AMBER)

add_text(s6, Inches(0.7), Inches(0.5), Inches(2), Inches(0.5),
         "06 / 06   CONCLUSION", size=14, bold=True, color=AMBER, font="Consolas")
add_speaker_chip(s6, "Personne 3")
add_text(s6, Inches(0.7), Inches(1.05), Inches(12), Inches(0.7),
         "Bilan et perspectives", size=32, bold=True, color=WHITE)

# left: coverage checklist
cov_x, cov_y = Inches(0.7), Inches(2.2)
cov_w, cov_h = Inches(6.0), Inches(4.4)
add_rounded(s6, cov_x, cov_y, cov_w, cov_h, NAVY)
add_text(s6, cov_x + Inches(0.35), cov_y + Inches(0.2), cov_w - Inches(0.55), Inches(0.4),
         "EXIGENCES COUVERTES", size=11, bold=True, color=AMBER, font="Consolas")

cov = [
    "GPIO  -  LEDs LD2, L0, L1 + bouton B1",
    "TIMER  -  TIM6 en mode interruption",
    "ADC  -  optionnel, retire pour simplifier",
    "UART  -  printf + saisie ville",
    "I2C  -  trois capteurs IKS01A3",
    "IRQ  -  TIM6 + EXTI13 + USART2 RX (3 sur 2)",
    "Capteurs  -  HTS221 + LPS22HH + STTS751 (3 sur 2)",
]
cy = cov_y + Inches(0.7)
for line in cov:
    # green dot
    dot = s6.shapes.add_shape(MSO_SHAPE.OVAL,
                              cov_x + Inches(0.35), cy + Inches(0.08),
                              Inches(0.22), Inches(0.22))
    dot.fill.solid(); dot.fill.fore_color.rgb = EMERALD
    dot.line.fill.background(); dot.shadow.inherit = False
    add_text(s6, cov_x + Inches(0.7), cy, cov_w - Inches(0.9), Inches(0.4),
             line, size=13, color=WHITE)
    cy += Inches(0.5)

# right: improvements + thanks
imp_x = Inches(7.0)
imp_w = Inches(5.6)
add_rounded(s6, imp_x, Inches(2.2), imp_w, Inches(2.3), TEAL)
add_text(s6, imp_x + Inches(0.35), Inches(2.4), imp_w - Inches(0.55), Inches(0.4),
         "PISTES D'AMELIORATION", size=11, bold=True, color=AMBER, font="Consolas")
imps = [
    "Affichage local sur MAX7219 (SPI1)",
    "Lecture ADC en DMA + journalisation EEPROM",
    "Calcul point de rosee + pression mer",
    "Transport sans fil : LoRa ou BLE",
]
iy = Inches(2.85)
for line in imps:
    add_text(s6, imp_x + Inches(0.35), iy, Inches(0.3), Inches(0.3),
             "►", size=12, color=AMBER, font="Consolas")
    add_text(s6, imp_x + Inches(0.7), iy, imp_w - Inches(0.9), Inches(0.3),
             line, size=13, color=WHITE)
    iy += Inches(0.35)

# Thank-you card
add_rounded(s6, imp_x, Inches(4.7), imp_w, Inches(1.9), AMBER)
add_text(s6, imp_x, Inches(4.85), imp_w, Inches(0.6),
         "MERCI POUR VOTRE ATTENTION", size=20, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s6, imp_x, Inches(5.55), imp_w, Inches(0.5),
         "Questions ?", size=18, italic=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s6, imp_x, Inches(6.05), imp_w, Inches(0.4),
         "Equipe : Personne 1 - Personne 2 - Personne 3", size=12,
         color=RGBColor(0xFE, 0xF3, 0xC7), align=PP_ALIGN.CENTER)

add_text(s6, Inches(0.7), Inches(7.05), Inches(11), Inches(0.35),
         "STM32 Meteo Station - NUCLEO-L152RE + X-NUCLEO-IKS01A3 - P23MRS 2026",
         size=10, color=RGBColor(0xCB, 0xD5, 0xE1), italic=True)

set_notes(s6, """[Personne 3] Pour conclure, on a couvert l'integralite du cahier des charges :

- les six peripheriques cles sont presents : GPIO, TIMER, UART et I2C. L'ADC etait optionnel, on l'a retire pour garder une demo plus lisible.
- on a trois peripheriques en interruption (TIM6, EXTI13, USART2 RX) alors que le sujet en demandait deux.
- on a trois capteurs actifs au lieu des deux exiges.

Cote pistes d'amelioration : on pourrait afficher la mesure localement sur le MAX7219, lire l'ADC en DMA pour faire de l'echantillonnage rapide, calculer des indicateurs derives comme le point de rosee, ou ajouter un transport sans fil - LoRa ou BLE - pour sortir le projet de la salle.

On retient surtout trois lecons : ne jamais bloquer dans une ISR, toujours verifier les retours BSP, et separer clairement la couche capteur de la couche application. Ca rend le code propre et reutilisable.

Merci pour votre attention - on est prets pour les questions.
""")


# ---------------------------------------------------------------------------
out = "STM32_METEO_STATION_presentation.pptx"
prs.save(out)
print("OK:", out)
