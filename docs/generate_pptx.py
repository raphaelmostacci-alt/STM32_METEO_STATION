"""
generate_pptx.py
----------------
Génère la présentation PowerPoint pour le projet STM32 Météo Station.
Projet : ISEN Toulon – P23MRS – Juin 2026

Dépendance : pip install python-pptx

Usage :
    python generate_pptx.py
    -> crée STM32_METEO_STATION_Presentation.pptx dans le même dossier
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Palette de couleurs (inspirée ST Microelectronics + ISEN)
# ---------------------------------------------------------------------------
ST_DARK   = RGBColor(0x03, 0x23, 0x4B)   # bleu marine ST
ST_CYAN   = RGBColor(0x00, 0xB4, 0xE3)   # cyan ST
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xF8)   # fond clair slides contenu
DARK_TXT  = RGBColor(0x1A, 0x1A, 0x2E)
MID_GREY  = RGBColor(0x6E, 0x7F, 0x8D)
GREEN_OK  = RGBColor(0x00, 0xA6, 0x51)
ORANGE    = RGBColor(0xFF, 0x6B, 0x00)

# ---------------------------------------------------------------------------
# Dimensions (16:9)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _solid_fill(shape, color: RGBColor):
    """Remplissage solide d'une forme."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def _no_fill(shape):
    shape.fill.background()


def _set_font(run, bold=False, italic=False, size_pt=None, color: RGBColor = None):
    run.font.bold = bold
    run.font.italic = italic
    if size_pt:
        run.font.size = Pt(size_pt)
    if color:
        run.font.color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor, radius=False):
    """Ajoute un rectangle coloré."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    _solid_fill(shape, color)
    shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, italic=False,
                color: RGBColor = WHITE, align=PP_ALIGN.LEFT,
                wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    _set_font(run, bold=bold, italic=italic, size_pt=font_size, color=color)
    return txBox


def add_bullet_textbox(slide, items: list, left, top, width, height,
                       title=None, font_size=16, title_size=20,
                       text_color: RGBColor = DARK_TXT,
                       title_color: RGBColor = ST_DARK):
    """Ajoute une zone de texte avec titre + bullets."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        _set_font(run, bold=True, size_pt=title_size, color=title_color)

    for item in items:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 1 if isinstance(item, tuple) and item[0] > 0 else 0
        text = item[1] if isinstance(item, tuple) else item
        run = p.add_run()
        run.text = ("    • " if (isinstance(item, tuple) and item[0] > 0) else "• ") + text
        _set_font(run, size_pt=font_size - 2 if (isinstance(item, tuple) and item[0] > 0) else font_size,
                  color=text_color)
    return txBox


def add_tag_chip(slide, text, left, top, width=Inches(1.8), height=Inches(0.35),
                 bg: RGBColor = ST_CYAN, fg: RGBColor = WHITE, size_pt=12):
    """Petit badge coloré (chip) pour les mots-clés."""
    chip = add_rect(slide, left, top, width, height, bg)
    tf = chip.text_frame
    tf.word_wrap = False
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    _set_font(run, bold=True, size_pt=size_pt, color=fg)
    return chip


def horizontal_line(slide, left, top, width, color: RGBColor = ST_CYAN, thickness=3):
    line = slide.shapes.add_shape(1, left, top, width, Pt(thickness))
    _solid_fill(line, color)
    line.line.fill.background()
    return line


# ---------------------------------------------------------------------------
# Layout commun pour slides de contenu
# ---------------------------------------------------------------------------

def _content_slide_base(prs, title_text, subtitle_text=None):
    """
    Fond blanc cassé, bandeau bleu en haut, titre blanc, numéro de slide.
    Retourne (slide, content_top, content_left, content_width, content_height).
    """
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Bandeau supérieur
    banner = add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), ST_DARK)

    # Bande cyan décorative fine
    horizontal_line(slide, 0, Inches(1.2), SLIDE_W, ST_CYAN, thickness=6)

    # Fond du corps (légèrement grisé)
    body_bg = add_rect(slide, 0, Inches(1.26), SLIDE_W, SLIDE_H - Inches(1.26), LIGHT_BG)

    # Titre dans le bandeau
    add_textbox(slide, title_text,
                Inches(0.4), Inches(0.15), Inches(11.0), Inches(0.85),
                font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    if subtitle_text:
        add_textbox(slide, subtitle_text,
                    Inches(0.4), Inches(0.78), Inches(9.0), Inches(0.38),
                    font_size=14, italic=True, color=ST_CYAN, align=PP_ALIGN.LEFT)

    # Petit logo textuel en haut à droite
    add_textbox(slide, "STM32 Météo Station",
                Inches(10.0), Inches(0.08), Inches(3.2), Inches(0.4),
                font_size=10, color=ST_CYAN, align=PP_ALIGN.RIGHT)

    return slide, Inches(1.35), Inches(0.35), SLIDE_W - Inches(0.7), SLIDE_H - Inches(1.55)


# ---------------------------------------------------------------------------
# SLIDE 1 – Titre
# ---------------------------------------------------------------------------

def slide_titre(prs):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Fond sombre total
    bg = add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, ST_DARK)

    # Bande déco gauche (cyan)
    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, ST_CYAN)

    # Bande déco bas
    add_rect(slide, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), ST_CYAN)

    # Zone centrale avec fond légèrement plus clair
    add_rect(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.5),
             RGBColor(0x06, 0x33, 0x6B))

    # Titre principal
    add_textbox(slide, "STM32 Météo Station",
                Inches(0.7), Inches(1.85), Inches(11.5), Inches(1.1),
                font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Ligne décorative cyan sous le titre
    horizontal_line(slide, Inches(2.5), Inches(3.0), Inches(8.3), ST_CYAN, thickness=4)

    # Sous-titre
    add_textbox(slide, "Station météorologique embarquée",
                Inches(0.7), Inches(3.1), Inches(11.5), Inches(0.6),
                font_size=22, italic=True, color=ST_CYAN, align=PP_ALIGN.CENTER)

    add_textbox(slide, "NUCLEO-L152RE  +  X-NUCLEO-IKS01A3",
                Inches(0.7), Inches(3.65), Inches(11.5), Inches(0.5),
                font_size=16, color=RGBColor(0xB0, 0xC8, 0xE0), align=PP_ALIGN.CENTER)

    # Membres
    add_textbox(slide, "Baptiste GREHL   ·   Mathias MOHA   ·   Raphael MOSTACCI",
                Inches(0.7), Inches(4.5), Inches(11.5), Inches(0.45),
                font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Date et école
    add_textbox(slide, "Projet STM32  –  ISEN Toulon  –  1er juin 2026",
                Inches(0.7), Inches(4.95), Inches(11.5), Inches(0.4),
                font_size=13, color=MID_GREY, align=PP_ALIGN.CENTER)

    # Chips technos
    chip_data = [
        ("C / HAL", ST_CYAN, DARK_TXT),
        ("STM32CubeIDE", ST_DARK, WHITE),
        ("I2C · UART · TIM · ADC", ORANGE, WHITE),
        ("3 capteurs shield", GREEN_OK, WHITE),
    ]
    chip_w = Inches(2.15)
    chip_h = Inches(0.33)
    gap = Inches(0.18)
    total_chips = len(chip_data)
    total_w = total_chips * chip_w + (total_chips - 1) * gap
    start_x = (SLIDE_W - total_w) / 2

    for i, (label, bg_c, fg_c) in enumerate(chip_data):
        x = start_x + i * (chip_w + gap)
        add_tag_chip(slide, label, x, Inches(5.6), chip_w, chip_h, bg_c, fg_c, size_pt=11)

    return slide


# ---------------------------------------------------------------------------
# SLIDE 2 – Présentation de l'application
# ---------------------------------------------------------------------------

def slide_presentation_app(prs):
    slide, ct, cl, cw, ch = _content_slide_base(
        prs, "Présentation de l'application",
        subtitle_text="Qu'est-ce que la STM32 Météo Station ?"
    )

    # Colonne gauche : description principale
    col_w = Inches(5.8)
    col_h = Inches(4.8)

    add_bullet_textbox(
        slide,
        [
            "Acquisition temps réel de données environnementales",
            (1, "Température (°C) – capteur STTS751 / HTS221"),
            (1, "Humidité relative (%) – capteur HTS221"),
            (1, "Pression atmosphérique (hPa) – capteur LPS22HH"),
            "Affichage des mesures via terminal série (UART)",
            (1, "Format : METEO Ville=X T=23.45 C RH=48.2 % P=1013.4 hPa"),
            "Deux modes de fonctionnement",
            (1, "LIVE : mesure et affichage toutes les secondes"),
            (1, "FROZEN : pause des mesures, LED mode allumée"),
            "Seuil d'alarme configurable par potentiomètre (ADC)",
            "Saisie de la ville d'observation via terminal UART",
        ],
        cl, ct, col_w, col_h,
        title="Fonctionnalités principales",
        font_size=14, title_size=18,
    )

    # Séparateur vertical
    add_rect(slide, cl + col_w + Inches(0.15), ct + Inches(0.1),
             Inches(0.04), Inches(4.5), ST_CYAN)

    # Colonne droite : schéma fonctionnel simplifié
    right_x = cl + col_w + Inches(0.4)
    right_w = SLIDE_W - right_x - Inches(0.3)

    add_textbox(slide, "Schéma fonctionnel",
                right_x, ct, right_w, Inches(0.4),
                font_size=18, bold=True, color=ST_DARK, align=PP_ALIGN.CENTER)

    blocks = [
        ("Shield IKS01A3", ST_CYAN, DARK_TXT, "HTS221 · LPS22HH · STTS751"),
        ("ADC1  PA0", ST_DARK, WHITE, "Potentiomètre → seuil alarme"),
        ("EXTI13  PC13", ORANGE, WHITE, "Bouton B1 → LIVE / FROZEN"),
        ("TIM6  1 Hz  IRQ", GREEN_OK, WHITE, "Tick périodique de mesure"),
        ("USART2  115200", RGBColor(0x6B, 0x48, 0xFF), WHITE, "Télémétrie + saisie ville"),
    ]

    bx = right_x + Inches(0.1)
    by = ct + Inches(0.55)
    bw = right_w - Inches(0.2)
    bh = Inches(0.68)
    gap = Inches(0.1)

    for (label, bg_c, fg_c, detail) in blocks:
        block = add_rect(slide, bx, by, bw, bh, bg_c)
        tf = block.text_frame
        tf.word_wrap = False
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = label
        _set_font(r1, bold=True, size_pt=13, color=RGBColor(*[int(c) for c in fg_c]))
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = detail
        _set_font(r2, size_pt=10, color=RGBColor(*[int(c) for c in fg_c]))
        by += bh + gap

    return slide


# ---------------------------------------------------------------------------
# SLIDE 3 – Solution Technique argumentée
# ---------------------------------------------------------------------------

def slide_solution_technique(prs):
    slide, ct, cl, cw, ch = _content_slide_base(
        prs, "Solution Technique",
        subtitle_text="Architecture matérielle et choix d'implémentation"
    )

    col_w = Inches(5.8)

    # Colonne gauche : périphériques
    add_bullet_textbox(
        slide,
        [
            "I2C1 (PB8/PB9) → shield IKS01A3",
            (1, "BSP IKS01A3 : abstraction multi-capteurs"),
            (1, "Capteurs : HTS221, LPS22HH, STTS751"),
            "USART2 (PA2/PA3) → printf retargeté",
            (1, "Réception ville par interruption RX"),
            "TIM6 Basic Timer → IRQ 1 Hz",
            (1, "Prescaler 31999, Period 999 @ 32 MHz"),
            "ADC1 IN0 (PA0) → 12 bits, polling",
            (1, "Conversion potentiomètre = seuil alarme"),
            "GPIO sorties → 5 LEDs de statut",
            (1, "LD2 heartbeat, L0 valid, L1 err, L2 mode, L3 alarm"),
        ],
        cl, ct, col_w, Inches(5.2),
        title="Périphériques utilisés",
        font_size=13, title_size=18,
    )

    # Séparateur
    add_rect(slide, cl + col_w + Inches(0.15), ct + Inches(0.1),
             Inches(0.04), Inches(4.6), ST_CYAN)

    # Colonne droite : interruptions + architecture logicielle
    right_x = cl + col_w + Inches(0.4)
    right_w = SLIDE_W - right_x - Inches(0.3)

    add_bullet_textbox(
        slide,
        [
            "2 périphériques en mode interruption ✓",
            (1, "TIM6 IRQ → flag g_sample_flag"),
            (1, "EXTI13 IRQ → flag g_ask_city + debounce 50 ms"),
            "Flags volatile partagés ISR ↔ main loop",
            (1, "g_sample_flag, g_paused, g_city_ready"),
            "Architecture non-bloquante",
            (1, "Pas de HAL_Delay dans la boucle principale"),
            (1, "UART RX par interruption (HAL_UART_Receive_IT)"),
            "Gestion d'erreurs capteurs",
            (1, "Indicateur valid dans WeatherSample"),
            (1, "LED erreur si init ou lecture KO"),
        ],
        right_x, ct, right_w, Inches(5.2),
        title="Architecture logicielle",
        font_size=13, title_size=18,
    )

    return slide


# ---------------------------------------------------------------------------
# SLIDE 4 – Défi & Problématique
# ---------------------------------------------------------------------------

def slide_defis(prs):
    slide, ct, cl, cw, ch = _content_slide_base(
        prs, "Défis & Problématiques",
        subtitle_text="Obstacles rencontrés pendant le développement"
    )

    challenges = [
        {
            "title": "Concurrence ISR / main loop",
            "color": ORANGE,
            "detail": (
                "TIM6 et EXTI13 peuvent lever des interruptions simultanément. "
                "Accès concurrent aux variables globales partagées."
            ),
        },
        {
            "title": "Saisie UART non-bloquante",
            "color": ST_CYAN,
            "detail": (
                "La saisie de la ville doit coexister avec le sampling 1 Hz "
                "sans bloquer la boucle principale ni perdre de caractères."
            ),
        },
        {
            "title": "Anti-rebond bouton (EXTI13)",
            "color": RGBColor(0x6B, 0x48, 0xFF),
            "detail": (
                "Le bouton bleu Nucleo génère plusieurs flancs parasites "
                "à chaque appui physique, causant des changements de mode involontaires."
            ),
        },
        {
            "title": "Initialisation conditionnelle des capteurs",
            "color": GREEN_OK,
            "detail": (
                "Le shield IKS01A3 peut ne pas avoir tous les capteurs soudés. "
                "L'initialisation partielle ne doit pas bloquer le démarrage."
            ),
        },
    ]

    card_w = Inches(5.9)
    card_h = Inches(1.85)
    gap_x = Inches(0.4)
    gap_y = Inches(0.2)

    positions = [
        (cl, ct),
        (cl + card_w + gap_x, ct),
        (cl, ct + card_h + gap_y),
        (cl + card_w + gap_x, ct + card_h + gap_y),
    ]

    for i, ch_data in enumerate(challenges):
        x, y = positions[i]
        # Carte fond blanc
        card = add_rect(slide, x, y, card_w, card_h, WHITE)
        card.line.color.rgb = ch_data["color"]
        card.line.width = Pt(2)

        # Barre de couleur gauche
        add_rect(slide, x, y, Inches(0.12), card_h, ch_data["color"])

        # Titre du défi
        add_textbox(slide, ch_data["title"],
                    x + Inches(0.2), y + Inches(0.1), card_w - Inches(0.3), Inches(0.42),
                    font_size=15, bold=True, color=ch_data["color"])

        # Détail
        add_textbox(slide, ch_data["detail"],
                    x + Inches(0.2), y + Inches(0.52), card_w - Inches(0.35), Inches(1.2),
                    font_size=12, color=DARK_TXT, wrap=True)

    return slide


# ---------------------------------------------------------------------------
# SLIDE 5 – Résolution
# ---------------------------------------------------------------------------

def slide_resolution(prs):
    slide, ct, cl, cw, ch = _content_slide_base(
        prs, "Résolution",
        subtitle_text="Solutions techniques mises en œuvre"
    )

    solutions = [
        {
            "problem": "Concurrence ISR / main",
            "solution": "Flags volatile uint8_t",
            "detail": "g_sample_flag, g_ask_city, g_city_ready, g_paused déclarés volatile – "
                      "la main loop lit et remet à zéro de façon atomique sur Cortex-M3.",
            "color": ORANGE,
        },
        {
            "problem": "UART non-bloquant",
            "solution": "HAL_UART_Receive_IT + callback",
            "detail": "Réception caractère par caractère en IRQ. HAL_UART_RxCpltCallback "
                      "accumule dans g_rx_buf et lève g_city_ready sur '\\n'.",
            "color": ST_CYAN,
        },
        {
            "problem": "Anti-rebond EXTI13",
            "solution": "Debounce 50 ms software",
            "detail": "HAL_GetTick() comparé à g_last_btn_tick dans le callback EXTI. "
                      "Si delta < 50 ms, l'événement est ignoré.",
            "color": RGBColor(0x6B, 0x48, 0xFF),
        },
        {
            "problem": "Capteurs partiels",
            "solution": "Init séquentielle indépendante",
            "detail": "Weather_InitSensors() tente chaque capteur séparément. "
                      "Flags s_hum_ok, s_press_ok, s_temp_ok mémorisent la disponibilité.",
            "color": GREEN_OK,
        },
    ]

    # Timeline-like layout : 4 étapes horizontales
    step_w = Inches(2.9)
    step_h = Inches(4.2)
    gap = Inches(0.18)
    start_x = cl

    for i, sol in enumerate(solutions):
        x = start_x + i * (step_w + gap)

        # Bloc principal
        block = add_rect(slide, x, ct, step_w, step_h, WHITE)
        block.line.color.rgb = sol["color"]
        block.line.width = Pt(1.5)

        # En-tête coloré
        header = add_rect(slide, x, ct, step_w, Inches(0.7), sol["color"])

        add_textbox(slide, f"#{i+1}",
                    x + Inches(0.08), ct + Inches(0.05), Inches(0.4), Inches(0.55),
                    font_size=22, bold=True, color=WHITE)

        add_textbox(slide, sol["solution"],
                    x + Inches(0.45), ct + Inches(0.08), step_w - Inches(0.55), Inches(0.58),
                    font_size=11, bold=True, color=WHITE, wrap=True)

        # Problème en petit
        add_textbox(slide, f"Problème : {sol['problem']}",
                    x + Inches(0.1), ct + Inches(0.75), step_w - Inches(0.2), Inches(0.4),
                    font_size=10, italic=True, color=MID_GREY, wrap=True)

        # Ligne séparatrice
        horizontal_line(slide, x + Inches(0.1), ct + Inches(1.15),
                        step_w - Inches(0.2), sol["color"], thickness=1)

        # Détail
        add_textbox(slide, sol["detail"],
                    x + Inches(0.1), ct + Inches(1.25), step_w - Inches(0.2), Inches(2.85),
                    font_size=11, color=DARK_TXT, wrap=True)

    return slide


# ---------------------------------------------------------------------------
# SLIDE 6 – Conclusion
# ---------------------------------------------------------------------------

def slide_conclusion(prs):
    slide, ct, cl, cw, ch = _content_slide_base(
        prs, "Conclusion",
        subtitle_text="Bilan et perspectives"
    )

    # Bilan – colonne gauche
    col_w = Inches(5.7)

    add_bullet_textbox(
        slide,
        [
            "Tous les critères du cahier des charges sont respectés",
            (1, "≥ 2 périphériques en mode interruption : TIM6 + EXTI13 ✓"),
            (1, "≥ 2 capteurs du shield : HTS221, LPS22HH, STTS751 ✓"),
            (1, "GPIO, TIMER, ADC, UART, I2C tous implémentés ✓"),
            "Code source propre et modulaire",
            (1, "Module weather.c / weather.h réutilisable"),
            (1, "Callbacks HAL respectés (pas de polling bloquant)"),
            "README.md documenté : branchement, build, utilisation",
            "Démonstration live : terminal série + LEDs de statut",
        ],
        cl, ct, col_w, Inches(4.5),
        title="Bilan du projet",
        font_size=14, title_size=18,
    )

    # Séparateur
    add_rect(slide, cl + col_w + Inches(0.15), ct + Inches(0.1),
             Inches(0.04), Inches(4.5), ST_CYAN)

    # Perspectives – colonne droite
    right_x = cl + col_w + Inches(0.4)
    right_w = SLIDE_W - right_x - Inches(0.3)

    add_bullet_textbox(
        slide,
        [
            "Améliorations possibles",
            (1, "Affichage LCD (SPI/I2C) pour autonomie"),
            (1, "Enregistrement SD card (SPI)"),
            (1, "Connectivité WiFi (ESP8266 UART)"),
            (1, "Horodatage RTC interne STM32L1"),
            "Optimisations",
            (1, "Mode Low-Power entre deux mesures"),
            (1, "DMA pour UART TX sans CPU"),
            "Enseignements",
            (1, "Maîtrise HAL + BSP STM32Cube"),
            (1, "Programmation événementielle / IRQ"),
            (1, "Debugging embarqué (SWD + printf)"),
        ],
        right_x, ct, right_w, Inches(4.5),
        title="Perspectives",
        font_size=13, title_size=18,
    )

    # Bande de bas de page – équipe
    add_rect(slide, 0, SLIDE_H - Inches(0.65), SLIDE_W, Inches(0.65), ST_DARK)
    add_textbox(slide,
                "Baptiste GREHL  ·  Mathias MOHA  ·  Raphael MOSTACCI  —  "
                "ISEN Toulon  –  Projet STM32  –  Juin 2026",
                Inches(0.3), SLIDE_H - Inches(0.58), SLIDE_W - Inches(0.6), Inches(0.5),
                font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

    return slide


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_titre(prs)
    slide_presentation_app(prs)
    slide_solution_technique(prs)
    slide_defis(prs)
    slide_resolution(prs)
    slide_conclusion(prs)

    import ctypes.wintypes, ctypes
    buf = ctypes.create_unicode_buffer(ctypes.wintypes.MAX_PATH)
    ctypes.windll.shell32.SHGetFolderPathW(None, 0x0000, None, 0, buf)
    out_dir  = buf.value  # chemin bureau natif (OneDrive inclus si sync activé)
    out_path = os.path.join(out_dir, "STM32_METEO_STATION_Presentation.pptx")
    prs.save(out_path)
    print(f"✅  Présentation générée : {out_path}")
    print(f"    {len(prs.slides)} slides  |  format 16:9 ({SLIDE_W/914400:.2f}\" × {SLIDE_H/914400:.2f}\")")


if __name__ == "__main__":
    build_presentation()
